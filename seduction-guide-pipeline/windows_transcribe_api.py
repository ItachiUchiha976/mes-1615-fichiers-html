#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
TRANSCRIPTION + EXTRACTION COMPLETE - API OPENAI WHISPER - WINDOWS 10
=======================================================================
Scanne TOUS les fichiers d'un dossier, affiche ce qu'il trouve,
puis traite automatiquement tout ce qui est utile :

  VIDEOS  (.mp4 .mkv .avi .mov...)  -> transcription API OpenAI Whisper
  AUDIOS  (.mp3 .wav .m4a .ogg...)  -> transcription API OpenAI Whisper
  PDFs    (.pdf)                    -> extraction texte (PyMuPDF, gratuit)
  WORD    (.docx)                   -> extraction texte (python-docx, gratuit)
  POWERPOINT (.pptx)               -> extraction texte des slides (python-pptx)
  TEXTE   (.txt .md .rtf)          -> lecture directe
  IMAGES  (.jpg .png .gif .webp...) -> liste des noms (pour reference visuelle)
  AUTRES  (.zip .exe .html...)      -> ignores (liste affichee)

COUT : $0.006 par minute audio/video uniquement
  Exemple : 50h video = ~$18

INSTALLATION :
  pip install openai pymupdf python-docx python-pptx

USAGE :
  1. Remplis OPENAI_API_KEY ci-dessous
  2. Double-clique sur ce fichier
  3. Selectionne ton dossier FORMATIONS
  4. Le script affiche tout ce qu'il a trouve et le cout estime
  5. Tu confirmes -> tout est traite en ~1h
"""

import json
import logging
import re
import subprocess
import sys
import tempfile
import time
from pathlib import Path
from datetime import datetime


# ================================================================================
#  CONFIGURATION
# ================================================================================

OPENAI_API_KEY = "sk-METS_TA_CLE_ICI"

OUTPUT_DIR       = "transcriptions"
CHECKPOINT_FILE  = "transcription_api_checkpoint.json"
LOG_FILE         = "transcription_api.log"

MAX_BYTES_PAR_PARTIE = int(3.5 * 1024 * 1024)
MAX_CHUNK_MB         = 24
CHUNK_DURATION_S     = 600

# ── Categories de fichiers ────────────────────────────────────────────────────

VIDEO_EXT = {".mp4", ".avi", ".mov", ".mkv", ".webm", ".mpeg", ".3gp", ".m4v", ".flv", ".wmv"}
AUDIO_EXT = {".mp3", ".wav", ".aac", ".ogg", ".wma", ".m4a", ".opus", ".flac"}
PDF_EXT   = {".pdf"}
WORD_EXT  = {".docx", ".doc"}
PPT_EXT   = {".pptx", ".ppt"}
TEXT_EXT  = {".txt", ".md", ".markdown", ".rtf", ".csv"}
IMAGE_EXT = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".tiff", ".tif", ".svg"}

# Tout ce qu'on traite (pas les images ni les autres)
PROCESS_EXT = VIDEO_EXT | AUDIO_EXT | PDF_EXT | WORD_EXT | PPT_EXT | TEXT_EXT
# Les images : on les liste seulement (pas de contenu extractible sans vision API)
LIST_ONLY_EXT = IMAGE_EXT

# ================================================================================


def setup_logging():
    fmt = "%(asctime)s  %(message)s"
    logging.basicConfig(
        level=logging.INFO,
        format=fmt,
        handlers=[
            logging.FileHandler(LOG_FILE, encoding="utf-8"),
            logging.StreamHandler(sys.stdout),
        ],
    )

log = logging.getLogger(__name__)


# -- Selecteur de dossier --------------------------------------------------------

def demander_dossier() -> Path:
    try:
        import tkinter as tk
        from tkinter import filedialog
        root = tk.Tk()
        root.withdraw()
        root.attributes("-topmost", True)
        print("\nUne fenetre de selection de dossier va s'ouvrir...")
        print("Choisis le dossier FORMATIONS sur ta cle USB.\n")
        dossier = filedialog.askdirectory(title="Choisis ton dossier FORMATIONS")
        root.destroy()
        return Path(dossier) if dossier else Path(".")
    except Exception as e:
        print(f"Fenetre impossible ({e}). Dossier courant utilise.")
        return Path(".")


# -- Tri naturel -----------------------------------------------------------------

def natural_key(name: str) -> list:
    return [int(c) if c.isdigit() else c.lower() for c in re.split(r"(\d+)", name)]

def sorted_naturally(paths):
    return sorted(paths, key=lambda p: natural_key(p.name))


# -- Checkpoint ------------------------------------------------------------------

def load_checkpoint() -> dict:
    if Path(CHECKPOINT_FILE).exists():
        try:
            return json.loads(Path(CHECKPOINT_FILE).read_text(encoding="utf-8"))
        except Exception:
            return {}
    return {}

def save_checkpoint(data: dict):
    Path(CHECKPOINT_FILE).write_text(
        json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8"
    )


# -- Scanner d'extensions --------------------------------------------------------

def scanner_extensions(root: Path) -> dict:
    """
    Parcourt tout le dossier et compte les fichiers par extension.
    Retourne un dict {extension: nombre}, trie par nombre decroissant.
    """
    counts = {}
    for f in root.rglob("*"):
        if f.is_file():
            ext = f.suffix.lower() if f.suffix else "(sans extension)"
            counts[ext] = counts.get(ext, 0) + 1
    return dict(sorted(counts.items(), key=lambda x: -x[1]))


def categoriser_fichier(ext: str) -> str:
    """Retourne la categorie d'un fichier selon son extension."""
    if ext in VIDEO_EXT:   return "VIDEO (Whisper API)"
    if ext in AUDIO_EXT:   return "AUDIO (Whisper API)"
    if ext in PDF_EXT:     return "PDF (extraction gratuite)"
    if ext in WORD_EXT:    return "WORD (extraction gratuite)"
    if ext in PPT_EXT:     return "POWERPOINT (extraction gratuite)"
    if ext in TEXT_EXT:    return "TEXTE (lecture directe)"
    if ext in IMAGE_EXT:   return "IMAGE (liste uniquement)"
    return "IGNORE"


# -- Duree des fichiers ----------------------------------------------------------

def get_duration_seconds(media_path: Path) -> float:
    cmd = ["ffprobe", "-v", "quiet", "-print_format", "json",
           "-show_format", str(media_path)]
    startupinfo = None
    if sys.platform == "win32":
        startupinfo = subprocess.STARTUPINFO()
        startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
    try:
        r = subprocess.run(cmd, capture_output=True, text=True,
                           startupinfo=startupinfo, timeout=30)
        return float(json.loads(r.stdout)["format"]["duration"])
    except Exception:
        return 0.0


# -- Extraction audio ------------------------------------------------------------

def extract_audio(media_path: Path, audio_path: Path):
    cmd = ["ffmpeg", "-y", "-i", str(media_path),
           "-vn", "-ac", "1", "-ar", "16000", "-b:a", "64k", str(audio_path)]
    startupinfo = None
    if sys.platform == "win32":
        startupinfo = subprocess.STARTUPINFO()
        startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
    r = subprocess.run(cmd, capture_output=True, text=True, startupinfo=startupinfo)
    if r.returncode != 0:
        raise RuntimeError(f"ffmpeg erreur : {r.stderr[-300:]}")


def split_audio_chunks(audio_path: Path, out_dir: Path) -> list:
    out_dir.mkdir(parents=True, exist_ok=True)
    pattern = str(out_dir / "chunk_%04d.mp3")
    cmd = ["ffmpeg", "-y", "-i", str(audio_path), "-f", "segment",
           "-segment_time", str(CHUNK_DURATION_S), "-c", "copy", pattern]
    startupinfo = None
    if sys.platform == "win32":
        startupinfo = subprocess.STARTUPINFO()
        startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
    subprocess.run(cmd, capture_output=True, startupinfo=startupinfo)
    return sorted(out_dir.glob("chunk_*.mp3"))


# -- Extraction texte : PDF, Word, PowerPoint, texte brut ------------------------

def extract_pdf(pdf_path: Path) -> tuple[str, int]:
    """Retourne (texte, nb_pages)."""
    try:
        import fitz
        doc = fitz.open(str(pdf_path))
        pages = []
        for i, page in enumerate(doc, 1):
            t = page.get_text().strip()
            if t:
                pages.append(f"--- Page {i} ---\n{t}")
        doc.close()
        if not pages:
            return "(PDF sans texte extractible - probablement scanne en image)", 0
        return "\n\n".join(pages), len(pages)
    except ImportError:
        return "(ERREUR : pip install pymupdf)", 0
    except Exception as e:
        return f"(Erreur PDF : {e})", 0


def extract_docx(docx_path: Path) -> str:
    """Extrait le texte d'un fichier Word .docx."""
    try:
        from docx import Document
        doc = Document(str(docx_path))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        # Inclure aussi les tableaux
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        return "\n\n".join(paragraphs) if paragraphs else "(Document Word vide)"
    except ImportError:
        return "(ERREUR : pip install python-docx)"
    except Exception as e:
        return f"(Erreur Word : {e})"


def extract_pptx(pptx_path: Path) -> tuple[str, int]:
    """Extrait le texte de chaque slide d'un PowerPoint."""
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
        if not slides:
            return "(PowerPoint sans texte extractible)", 0
        return "\n\n".join(slides), len(slides)
    except ImportError:
        return "(ERREUR : pip install python-pptx)", 0
    except Exception as e:
        return f"(Erreur PowerPoint : {e})", 0


def extract_text_file(txt_path: Path) -> str:
    """Lit un fichier texte brut (.txt, .md, .csv, .rtf)."""
    for encoding in ("utf-8", "latin-1", "cp1252"):
        try:
            return txt_path.read_text(encoding=encoding).strip()
        except Exception:
            continue
    return "(Impossible de lire ce fichier texte)"


# -- Traitement d'un fichier non-media (PDF/Word/PPT/texte) ----------------------

def process_text_file(file_path: Path, output_txt: Path, checkpoint: dict) -> bool:
    key = str(file_path)
    if checkpoint.get(key) == "ok" and output_txt.exists():
        log.info(f"  [SKIP] {file_path.name}")
        return True

    ext = file_path.suffix.lower()
    output_txt.parent.mkdir(parents=True, exist_ok=True)

    try:
        if ext in PDF_EXT:
            text, nb = extract_pdf(file_path)
            type_str = f"PDF ({nb} pages)"
            log.info(f"  --> [PDF] {file_path.name} ({nb} pages)")
        elif ext in WORD_EXT:
            text = extract_docx(file_path)
            type_str = "WORD"
            log.info(f"  --> [WORD] {file_path.name}")
        elif ext in PPT_EXT:
            text, nb = extract_pptx(file_path)
            type_str = f"POWERPOINT ({nb} slides)"
            log.info(f"  --> [PPT] {file_path.name} ({nb} slides)")
        elif ext in TEXT_EXT:
            text = extract_text_file(file_path)
            type_str = "TEXTE"
            log.info(f"  --> [TXT] {file_path.name}")
        else:
            return False

        header = (
            f"=== FICHIER : {file_path.name} ===\n"
            f"=== TYPE : {type_str} ===\n"
            f"=== DOSSIER : {file_path.parent.name} ===\n\n"
        )
        output_txt.write_text(header + text + "\n", encoding="utf-8")
        checkpoint[key] = "ok"
        save_checkpoint(checkpoint)
        log.info(f"    OK : {output_txt.name}")
        return True
    except Exception as e:
        log.error(f"    ECHEC {file_path.name} : {e}")
        checkpoint[key] = f"erreur: {e}"
        save_checkpoint(checkpoint)
        return False


# -- Traitement images : copie + PDF regroupe ------------------------------------

def process_images(images: list, output_dir: Path, media_root: Path):
    """
    Pour chaque image trouvee :
      1. La copie dans output_dir/Images_utiles/ (meme sous-dossier que l'original)
      2. Cree un PDF unique 'TOUTES_LES_IMAGES.pdf' avec toutes les images,
         une par page, avec le nom du fichier et son dossier en titre de page.

    Le PDF peut etre envoye directement a Claude sur claude.ai (1 seul fichier).
    """
    import shutil

    if not images:
        return

    output_dir.mkdir(parents=True, exist_ok=True)
    images_dir = output_dir / "Images_utiles"
    images_dir.mkdir(exist_ok=True)

    log.info(f"  Copie de {len(images)} image(s) dans Images_utiles/...")

    # Copier les images en preservant l'arborescence
    copied = []
    for img in sorted_naturally(images):
        try:
            rel = img.relative_to(media_root)
        except ValueError:
            rel = Path(img.name)
        dest = images_dir / rel
        dest.parent.mkdir(parents=True, exist_ok=True)
        try:
            shutil.copy2(str(img), str(dest))
            copied.append((img, rel))
        except Exception as e:
            log.warning(f"    Copie impossible : {img.name} -> {e}")

    log.info(f"  {len(copied)} image(s) copiee(s) dans : {images_dir}")

    # Creer le PDF regroupe avec PyMuPDF
    pdf_output = output_dir / "TOUTES_LES_IMAGES.pdf"
    try:
        import fitz  # PyMuPDF

        doc = fitz.open()  # nouveau document PDF vide

        for img_path, rel in copied:
            try:
                # Ouvrir l'image et l'inserer dans une nouvelle page
                img_doc = fitz.open(str(img_path))
                # Recuperer les dimensions de l'image
                img_rect = img_doc[0].rect if img_doc.page_count > 0 else fitz.Rect(0, 0, 595, 842)
                img_doc.close()

                # Taille de page : A4 portrait ou dimensions de l'image si plus grande
                page_w = max(595, img_rect.width + 40)
                page_h = max(100, img_rect.height + 80)  # +80 pour le titre en haut

                page = doc.new_page(width=page_w, height=page_h)

                # Titre en haut de page : nom du fichier + dossier
                titre = f"{img_path.name}  |  {rel.parent}"
                page.insert_text(
                    fitz.Point(20, 20),
                    titre,
                    fontsize=10,
                    color=(0.2, 0.2, 0.2),
                )
                page.insert_text(
                    fitz.Point(20, 35),
                    "-" * min(80, int(page_w / 7)),
                    fontsize=8,
                    color=(0.5, 0.5, 0.5),
                )

                # Inserer l'image sous le titre
                img_zone = fitz.Rect(20, 50, page_w - 20, page_h - 20)
                page.insert_image(img_zone, filename=str(img_path))

            except Exception as e:
                # Si une image pose probleme, ajouter une page d'erreur
                page = doc.new_page(width=595, height=100)
                page.insert_text(
                    fitz.Point(20, 50),
                    f"ERREUR image : {img_path.name} -> {e}",
                    fontsize=10,
                    color=(0.8, 0, 0),
                )

        doc.save(str(pdf_output))
        doc.close()
        taille_mo = pdf_output.stat().st_size / 1024 / 1024
        log.info(f"  PDF images cree : {pdf_output}  ({taille_mo:.1f} Mo)")
        log.info(f"  -> Envoie ce PDF a Claude en meme temps que les transcriptions.")

    except ImportError:
        log.warning("  PyMuPDF non installe -> PDF non cree. Lance : pip install pymupdf")
    except Exception as e:
        log.error(f"  Erreur creation PDF images : {e}")

    # Creer aussi le fichier texte de reference
    ref_file = output_dir / "LISTE_IMAGES.txt"
    with open(ref_file, "w", encoding="utf-8") as f:
        f.write("=== IMAGES DES FORMATIONS ===\n")
        f.write(f"=== Total : {len(copied)} images ===\n\n")
        f.write("FICHIER PDF CONTENANT TOUTES LES IMAGES : TOUTES_LES_IMAGES.pdf\n")
        f.write("Envoie ce PDF a Claude en meme temps que les fichiers de transcription.\n\n")
        current_folder = None
        for img_path, rel in copied:
            folder = str(rel.parent)
            if folder != current_folder:
                f.write(f"\n[{folder}]\n")
                current_folder = folder
            f.write(f"  - {img_path.name}\n")
    log.info(f"  Liste de reference : {ref_file}")


# -- Transcription API OpenAI ----------------------------------------------------

def transcribe_with_api(audio_path: Path) -> tuple:
    from openai import OpenAI
    client = OpenAI(api_key=OPENAI_API_KEY)
    file_size_mb = audio_path.stat().st_size / (1024 * 1024)

    if file_size_mb <= MAX_CHUNK_MB:
        with open(audio_path, "rb") as f:
            resp = client.audio.transcriptions.create(
                model="whisper-1", file=f, response_format="verbose_json"
            )
        return resp.text.strip(), getattr(resp, "language", "en")
    else:
        log.info(f"    Fichier {file_size_mb:.0f} Mo -> decoupage en segments 10 min...")
        chunks_dir = audio_path.parent / (audio_path.stem + "_chunks")
        chunks = split_audio_chunks(audio_path, chunks_dir)
        texts, lang = [], "en"
        for j, chunk in enumerate(chunks, 1):
            log.info(f"    Segment {j}/{len(chunks)}...")
            with open(chunk, "rb") as f:
                resp = client.audio.transcriptions.create(
                    model="whisper-1", file=f, response_format="verbose_json"
                )
            texts.append(resp.text.strip())
            if j == 1:
                lang = getattr(resp, "language", "en")
        return " ".join(texts), lang


def process_media(media_path: Path, output_txt: Path, checkpoint: dict) -> tuple:
    key = str(media_path)
    if checkpoint.get(key) == "ok" and output_txt.exists():
        log.info(f"  [SKIP] {media_path.name}")
        return True, 0.0

    type_f = "AUDIO" if media_path.suffix.lower() in AUDIO_EXT else "VIDEO"
    log.info(f"  --> [{type_f}] {media_path.name}")
    output_txt.parent.mkdir(parents=True, exist_ok=True)
    t_start = time.time()

    for attempt in range(1, 4):
        try:
            with tempfile.TemporaryDirectory() as tmp:
                audio = Path(tmp) / "audio.mp3"
                extract_audio(media_path, audio)
                text, lang = transcribe_with_api(audio)
            # Verifier si un fichier de sous-titres existe deja pour cette video
            # (meme nom, extension .txt, dans le meme dossier)
            existing_subtitle = output_txt  # meme chemin cible
            subtitle_content = ""
            if existing_subtitle.exists():
                try:
                    existing = existing_subtitle.read_text(encoding="utf-8")
                    if "SOUS-TITRES" in existing or "sous-tit" in existing.lower():
                        subtitle_content = existing
                except Exception:
                    pass

            header = (
                f"=== FICHIER : {media_path.name} ===\n"
                f"=== TYPE : {type_f} ===\n"
                f"=== DOSSIER : {media_path.parent.name} ===\n"
                f"=== LANGUE DETECTEE : {lang} ===\n\n"
            )

            if subtitle_content:
                # Fusionner : sous-titres existants + transcription Whisper complete
                contenu_final = (
                    header +
                    "=== TRANSCRIPTION WHISPER (complete) ===\n\n" +
                    text + "\n\n" +
                    "=== SOUS-TITRES ORIGINAUX (partiels, pour reference) ===\n\n" +
                    subtitle_content
                )
            else:
                contenu_final = header + text + "\n"

            output_txt.write_text(contenu_final, encoding="utf-8")
            checkpoint[key] = "ok"
            save_checkpoint(checkpoint)
            duree = time.time() - t_start
            log.info(f"    OK ({lang}) en {format_duree(duree)}")
            return True, duree
        except Exception as e:
            if attempt < 3:
                wait = 15 * attempt
                log.warning(f"    Tentative {attempt}/3 : {e}. Pause {wait}s...")
                time.sleep(wait)
            else:
                log.error(f"    ECHEC : {media_path.name} -> {e}")
                checkpoint[key] = f"erreur: {e}"
                save_checkpoint(checkpoint)
                return False, 0.0


# -- Suivi de progression --------------------------------------------------------

def format_duree(secondes: float) -> str:
    s = int(secondes)
    h, m, sec = s // 3600, (s % 3600) // 60, s % 60
    if h > 0:
        return f"{h}h{m:02d}m{sec:02d}s"
    return f"{m}m{sec:02d}s"

def affiche_barre(i, total, debut, success, errors, durees):
    pct = i / total * 100
    elapsed = time.time() - debut
    eta = format_duree((sum(durees)/len(durees)) * (total-i)) if durees else "calcul..."
    largeur = 26
    rempli = int(largeur * i / total)
    barre = "#" * rempli + "-" * (largeur - rempli)
    log.info(
        f"  [{barre}] {pct:5.1f}%  {i}/{total}  |  "
        f"OK:{success} Err:{errors}  |  "
        f"Ecoule:{format_duree(elapsed)}  Reste:{eta}"
    )


# -- Fusion finale ---------------------------------------------------------------

def fusionne_tout(output_dir: Path, base_name: str = "transcription_COMPLETE"):
    log.info("\nFusion et decoupage en parties de 3.5 Mo...")

    dossiers = {}
    for txt in output_dir.rglob("*.txt"):
        dossiers.setdefault(txt.parent, []).append(txt)

    dossiers_tries = sorted(dossiers.keys(),
                            key=lambda p: (-len(p.parts), natural_key(p.name)))
    sections = []
    for dossier in dossiers_tries:
        txts = sorted_naturally(dossiers[dossier])
        try:
            rel = dossier.relative_to(output_dir)
            nom = str(rel) if str(rel) != "." else "Racine"
        except ValueError:
            nom = dossier.name
        sections.append((nom, txts))

    total_f = sum(len(t) for _, t in sections)
    date_str = datetime.now().strftime("%d/%m/%Y %H:%M")

    def ouvre_partie(num):
        nom_f = f"{base_name}_PARTIE_{num}.txt"
        f = open(nom_f, "w", encoding="utf-8")
        f.write("=" * 70 + "\n")
        f.write("  TRANSCRIPTIONS + DOCUMENTS - FORMATIONS SUR LA SEDUCTION\n")
        f.write(f"  PARTIE {num} | {date_str} | {total_f} fichiers au total\n")
        f.write("=" * 70 + "\n\n")
        return f, nom_f

    partie_num = 1
    f, nom_courant = ouvre_partie(partie_num)
    octets = 0
    fichiers_crees = [nom_courant]

    # Sommaire
    lignes = [
        "+" + "-" * 68 + "+\n",
        "|  SOMMAIRE - ARBORESCENCE COMPLETE DES FORMATIONS                |\n",
        "|  Videos, audios, PDFs, Word, PowerPoint, textes                 |\n",
        "+" + "-" * 68 + "+\n\n"
    ]
    num = 1
    for nom_section, txts in sections:
        depth = nom_section.count("\\") + nom_section.count("/")
        indent = "    " * depth
        lignes.append(f"{indent}[DOSSIER] {nom_section}\n")
        for txt in txts:
            lignes.append(f"{indent}    [{num:03d}] {txt.stem}\n")
            num += 1
        lignes.append("\n")
    lignes += ["\n" + "=" * 70 + "\n", "  FIN DU SOMMAIRE\n", "=" * 70 + "\n\n"]
    bloc = "".join(lignes)
    f.write(bloc)
    octets += len(bloc.encode("utf-8"))

    for nom_section, txts in sections:
        sh = ("\n" + "-" * 60 + "\n" + f"  SECTION : {nom_section}\n"
              + "-" * 60 + "\n\n")
        shb = len(sh.encode("utf-8"))
        if octets + shb > MAX_BYTES_PAR_PARTIE:
            f.close()
            partie_num += 1
            f, nom_courant = ouvre_partie(partie_num)
            fichiers_crees.append(nom_courant)
            octets = 0
        f.write(sh)
        octets += shb
        for txt in txts:
            contenu = txt.read_text(encoding="utf-8") + "\n\n"
            cb = len(contenu.encode("utf-8"))
            if octets > 0 and octets + cb > MAX_BYTES_PAR_PARTIE:
                f.close()
                partie_num += 1
                f, nom_courant = ouvre_partie(partie_num)
                fichiers_crees.append(nom_courant)
                octets = 0
            f.write(contenu)
            octets += cb
    f.close()

    log.info(f"Fusion terminee : {partie_num} partie(s)")
    for nf in fichiers_crees:
        taille = Path(nf).stat().st_size / 1024 / 1024
        log.info(f"  {nf}  ({taille:.1f} Mo)")
    return fichiers_crees


# -- Programme principal ---------------------------------------------------------

def main():
    setup_logging()

    # Verifier cle API
    if OPENAI_API_KEY == "sk-METS_TA_CLE_ICI":
        print("\n" + "=" * 70)
        print("  ERREUR : Cle API OpenAI manquante !")
        print("=" * 70)
        print("\nOuvre ce fichier avec Notepad++ et remplace :")
        print("  OPENAI_API_KEY = \"sk-METS_TA_CLE_ICI\"")
        print("par ta vraie cle.")
        print("\nObtenir une cle : https://platform.openai.com/api-keys")
        input("\nAppuie sur Entree pour fermer...")
        sys.exit(1)

    # Verifier ffmpeg
    try:
        subprocess.run(["ffmpeg", "-version"], capture_output=True, check=True)
    except (FileNotFoundError, subprocess.CalledProcessError):
        print("ERREUR : ffmpeg introuvable.")
        print("Lance : winget install ffmpeg  (cmd administrateur)")
        input("\nAppuie sur Entree pour fermer...")
        sys.exit(1)

    # Choisir le dossier
    media_root = Path(sys.argv[1]) if len(sys.argv) > 1 else demander_dossier()
    media_root = media_root.resolve()
    output_dir = Path(OUTPUT_DIR).resolve()

    log.info("=" * 70)
    log.info("  TRANSCRIPTION + EXTRACTION COMPLETE")
    log.info("=" * 70)
    log.info(f"  Dossier source : {media_root}")
    log.info(f"  Sorties .txt   : {output_dir}")
    log.info("=" * 70 + "\n")

    # ── SCAN DE TOUTES LES EXTENSIONS ────────────────────────────────────────
    log.info("Scan du dossier en cours...\n")
    counts = scanner_extensions(media_root)

    if not counts:
        log.error(f"Dossier vide ou introuvable : {media_root}")
        input("Appuie sur Entree pour fermer...")
        sys.exit(1)

    # Afficher le tableau complet
    log.info("=" * 70)
    log.info("  FICHIERS TROUVES PAR TYPE")
    log.info("=" * 70)
    log.info(f"  {'Extension':<15} {'Nombre':>8}    {'Action'}")
    log.info(f"  {'-'*15} {'-'*8}    {'-'*35}")

    total_tous = 0
    for ext, nb in counts.items():
        categorie = categoriser_fichier(ext)
        log.info(f"  {ext:<15} {nb:>8}    {categorie}")
        total_tous += nb

    log.info(f"  {'-'*15} {'-'*8}")
    log.info(f"  {'TOTAL':<15} {total_tous:>8}")
    log.info("=" * 70 + "\n")

    # Collecter les fichiers par categorie
    tous = sorted_naturally([p for p in media_root.rglob("*") if p.is_file()])

    medias   = [f for f in tous if f.suffix.lower() in VIDEO_EXT | AUDIO_EXT]
    docs     = [f for f in tous if f.suffix.lower() in PDF_EXT | WORD_EXT | PPT_EXT | TEXT_EXT]
    images   = [f for f in tous if f.suffix.lower() in IMAGE_EXT]
    ignores  = [f for f in tous if f.suffix.lower() not in PROCESS_EXT | LIST_ONLY_EXT]

    nb_v = sum(1 for f in medias if f.suffix.lower() in VIDEO_EXT)
    nb_a = sum(1 for f in medias if f.suffix.lower() in AUDIO_EXT)
    nb_p = sum(1 for f in docs   if f.suffix.lower() in PDF_EXT)
    nb_w = sum(1 for f in docs   if f.suffix.lower() in WORD_EXT)
    nb_x = sum(1 for f in docs   if f.suffix.lower() in PPT_EXT)
    nb_t = sum(1 for f in docs   if f.suffix.lower() in TEXT_EXT)

    log.info(f"RESUME DES TRAITEMENTS :")
    log.info(f"  {nb_v} videos         -> Whisper API (payant)")
    log.info(f"  {nb_a} audios         -> Whisper API (payant)")
    log.info(f"  {nb_p} PDFs           -> PyMuPDF (gratuit)")
    log.info(f"  {nb_w} Word (.docx)   -> python-docx (gratuit)")
    log.info(f"  {nb_x} PowerPoint     -> python-pptx (gratuit)")
    log.info(f"  {nb_t} fichiers texte -> lecture directe (gratuit)")
    log.info(f"  {len(images)} images       -> copie dans Images_utiles/ + TOUTES_LES_IMAGES.pdf")
    log.info(f"  {len(ignores)} autres       -> ignores")

    if not medias and not docs:
        log.error("Aucun fichier a traiter trouve.")
        input("Appuie sur Entree pour fermer...")
        sys.exit(1)

    # ── ESTIMATION DU COUT ───────────────────────────────────────────────────
    if medias:
        log.info(f"\nCalcul de la duree totale des {len(medias)} videos/audios...\n")
        duree_totale_s = 0.0
        for i, m in enumerate(medias, 1):
            duree_totale_s += get_duration_seconds(m)
            print(f"\r  Analyse {i}/{len(medias)} : {m.name[:50]:<50}", end="", flush=True)
        print()
    else:
        duree_totale_s = 0.0

    duree_min = duree_totale_s / 60
    cout = duree_min * 0.006
    marge = cout + 3

    log.info("\n" + "=" * 70)
    log.info("  ESTIMATION DU COUT")
    log.info("=" * 70)
    if medias:
        log.info(f"  Duree totale videos/audios : {format_duree(duree_totale_s)} ({duree_min:.0f} min)")
        log.info(f"  Tarif Whisper API          : $0.006 / minute")
        log.info(f"  COUT ESTIME                : ${cout:.2f}")
        log.info(f"  Credits recommandes        : au moins ${marge:.0f}")
    else:
        log.info("  Aucun video/audio -> cout API : $0.00")
    log.info(f"  Documents/PDFs/textes      : GRATUIT")
    log.info(f"  Duree totale estimee       : ~30 a 90 minutes")
    log.info("=" * 70)
    if medias:
        log.info(f"\n  Verifie tes credits : https://platform.openai.com/usage")
    log.info("")

    reponse = input("Tout traiter maintenant ? (oui / non) : ").strip().lower()
    if reponse not in ("oui", "o", "yes", "y"):
        print("Annule.")
        sys.exit(0)

    checkpoint = load_checkpoint()
    output_dir.mkdir(parents=True, exist_ok=True)

    # ── ETAPE 1 : Documents (gratuit) ────────────────────────────────────────
    if docs:
        log.info(f"\n--- ETAPE 1/3 : Documents ({len(docs)} fichiers, GRATUIT) ---\n")
        ok_d, err_d = 0, 0
        for i, doc in enumerate(docs, 1):
            try:
                rel = doc.relative_to(media_root)
            except ValueError:
                rel = Path(doc.name)
            output_txt = output_dir / rel.parent / (doc.stem + ".txt")
            log.info(f"[{i}/{len(docs)}] {rel}")
            if process_text_file(doc, output_txt, checkpoint):
                ok_d += 1
            else:
                err_d += 1
        log.info(f"\nDocuments : {ok_d} OK, {err_d} erreur(s)")

    # ── ETAPE 2 : Images (copie + PDF regroupe) ─────────────────────────────
    if images:
        log.info(f"\n--- ETAPE 2/3 : Images ({len(images)} -> copie + PDF regroupe) ---")
        process_images(images, output_dir, media_root)

    # ── ETAPE 3 : Videos/Audios (API Whisper) ────────────────────────────────
    if medias:
        total = len(medias)
        deja = sum(1 for m in medias if checkpoint.get(str(m)) == "ok")
        log.info(f"\n--- ETAPE 3/3 : Videos/Audios ({total} fichiers, API) ---")
        log.info(f"  Deja transcrits   : {deja} [SKIP]")
        log.info(f"  Restant a traiter : {total - deja}\n")

        success, errors, durees = 0, 0, []
        debut_global = time.time()

        for i, media in enumerate(medias, 1):
            try:
                rel = media.relative_to(media_root)
            except ValueError:
                rel = Path(media.name)
            output_txt = output_dir / rel.parent / (media.stem + ".txt")
            log.info(f"\n[{i}/{total}] {rel}")
            ok, duree = process_media(media, output_txt, checkpoint)
            if ok:
                success += 1
                if duree > 1.0:
                    durees.append(duree)
            else:
                errors += 1
            affiche_barre(i, total, debut_global, success, errors, durees)

        log.info("\n" + "=" * 70)
        log.info(f"  Videos/Audios : {success}/{total} OK  |  Erreurs : {errors}")
        log.info(f"  Cout reel     : ~${cout:.2f}")
        log.info("=" * 70)

    # ── Fusion finale ────────────────────────────────────────────────────────
    fichiers = fusionne_tout(output_dir)
    nb = len(fichiers)

    log.info(f"\n{'='*70}")
    log.info(f"  TOUT EST TERMINE !")
    log.info(f"  {nb} fichier(s) pret(s) pour Claude :")
    for fich in fichiers:
        log.info(f"  {Path(fich).resolve()}")
    # Verifier si le PDF images existe
    pdf_images = output_dir / "TOUTES_LES_IMAGES.pdf"
    if pdf_images.exists():
        taille_img = pdf_images.stat().st_size / 1024 / 1024
        log.info(f"\n  PDF images : {pdf_images}  ({taille_img:.1f} Mo)")

    log.info("\n  SUR CLAUDE.AI - ORDRE D'ENVOI :")
    if nb > 1:
        log.info("  1. 'Je vais te donner plusieurs fichiers de formations sur la")
        log.info("     seduction (transcriptions videos + PDFs + docs + images).")
        log.info("     Lis TOUT avant de creer le guide. Voici la partie 1.'")
        log.info("     [attache PARTIE_1.txt]")
        log.info("  2. 'Voici la suite.' [attache PARTIE_2.txt, PARTIE_3.txt...]")
    else:
        log.info("  1. 'Voici les transcriptions et documents de mes formations.'")
        log.info("     [attache le fichier PARTIE_1.txt]")
    if pdf_images.exists():
        log.info("  +  'Voici egalement les images/schemas des formations.'")
        log.info("     [attache TOUTES_LES_IMAGES.pdf]")
    log.info("  -> Dernier message : 'C'est tout. Cree maintenant le guide")
    log.info("     COMPLET de seduction en francais, en 9 chapitres :'")
    log.info("     1. Mindset & confiance en soi")
    log.info("     2. Comprendre les femmes")
    log.info("     3. L'approche (avec scripts)")
    log.info("     4. La conversation & connexion")
    log.info("     5. La seduction progressive")
    log.info("     6. Situations specifiques (Tinder, soirees, etc.)")
    log.info("     7. Erreurs a eviter")
    log.info("     8. Scripts & phrases pratiques")
    log.info("     9. Developpement personnel'")
    log.info(f"{'='*70}\n")

    input("Appuie sur Entree pour fermer...")


if __name__ == "__main__":
    main()
